from __future__ import annotations

from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "DynaTwin-Swarm路演PPT.pptx"

W, H = Inches(13.333333), Inches(7.5)

BG = RGBColor(4, 14, 36)
PANEL = RGBColor(12, 31, 65)
PANEL_2 = RGBColor(16, 46, 88)
CYAN = RGBColor(0, 212, 255)
WHITE = RGBColor(245, 250, 255)
MUTED = RGBColor(154, 181, 210)
DIM = RGBColor(94, 125, 158)
AMBER = RGBColor(255, 184, 0)
RED = RGBColor(255, 82, 82)
GREEN = RGBColor(0, 255, 136)
VIOLET = RGBColor(137, 102, 255)

FONT = "Microsoft YaHei"
MONO = "Consolas"


def cm(v: float):
    return Inches(v / 2.54)


def set_text(frame, text: str, size: int, color=WHITE, bold=False, align=None, font=FONT):
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return p


def add_text(slide, text, x, y, w, h, size=20, color=WHITE, bold=False, align=None, font=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    set_text(box.text_frame, text, size, color, bold, align, font)
    return box


def add_multiline(slide, lines, x, y, w, h, size=18, color=WHITE, bullet=False, gap=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        if bullet:
            p.level = 0
            p._p.get_or_add_pPr().set("marL", "250000")
    return box


def add_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    # soft cyan/violet halos
    for x, y, r, color, trans in [
        (Inches(-0.6), Inches(-0.45), Inches(3.3), CYAN, 74),
        (Inches(10.2), Inches(0.2), Inches(3.0), VIOLET, 80),
        (Inches(6.0), Inches(5.7), Inches(3.6), CYAN, 86),
    ]:
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, r, r)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = trans
        shp.line.fill.background()

    # grid
    for i in range(0, 16):
        x = Inches(0.4 + i * 0.85)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Inches(0.15), x, Inches(7.25))
        line.line.color.rgb = RGBColor(16, 55, 95)
        line.line.transparency = 72
        line.line.width = Pt(0.4)
    for i in range(0, 9):
        y = Inches(0.3 + i * 0.82)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), y, Inches(13.15), y)
        line.line.color.rgb = RGBColor(16, 55, 95)
        line.line.transparency = 78
        line.line.width = Pt(0.4)


def glass_rect(slide, x, y, w, h, radius=True, fill=PANEL, line=CYAN, transparency=22):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    shp.line.color.rgb = line
    shp.line.transparency = 44
    shp.line.width = Pt(1.1)
    return shp


def title(slide, heading, subtitle=None, num=None):
    add_text(slide, heading, Inches(0.7), Inches(0.52), Inches(8.4), Inches(0.55), 28, WHITE, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.72), Inches(1.06), Inches(10.2), Inches(0.32), 12, MUTED)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(1.32), Inches(12.65), Inches(1.32))
    line.line.color.rgb = CYAN
    line.line.transparency = 22
    line.line.width = Pt(1.2)
    if num:
        add_text(slide, f"{num:02d}", Inches(11.8), Inches(0.45), Inches(0.8), Inches(0.42), 18, CYAN, True, PP_ALIGN.RIGHT, MONO)


def add_tag(slide, text, x, y, w=1.2, color=CYAN):
    shp = glass_rect(slide, x, y, w, Inches(0.32), fill=RGBColor(5, 30, 58), line=color, transparency=8)
    set_text(shp.text_frame, text, 10, color, True, PP_ALIGN.CENTER)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def card(slide, x, y, w, h, head, body, icon=None, accent=CYAN):
    shp = glass_rect(slide, x, y, w, h)
    if icon:
        add_text(slide, icon, x + Inches(0.18), y + Inches(0.17), Inches(0.42), Inches(0.42), 22, accent, True)
        tx = x + Inches(0.68)
        tw = w - Inches(0.86)
    else:
        tx = x + Inches(0.24)
        tw = w - Inches(0.48)
    add_text(slide, head, tx, y + Inches(0.22), tw, Inches(0.42), 18, WHITE, True)
    add_multiline(slide, body.split("\n"), tx, y + Inches(0.78), tw, h - Inches(0.9), 13, MUTED, gap=2)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.05), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    return shp


def tech_decoration(slide):
    cx, cy = Inches(10.6), Inches(3.6)
    for r, trans in [(2.4, 64), (1.75, 72), (1.12, 78)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(r / 2), cy - Inches(r / 2), Inches(r), Inches(r))
        c.fill.background()
        c.line.color.rgb = CYAN
        c.line.transparency = trans
        c.line.width = Pt(1.2)
    for a in range(0, 360, 45):
        x2 = cx + Inches(math.cos(math.radians(a)) * 1.55)
        y2 = cy + Inches(math.sin(math.radians(a)) * 1.55)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, x2, y2)
        line.line.color.rgb = CYAN
        line.line.transparency = 42
        line.line.width = Pt(1.1)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x2 - Inches(0.06), y2 - Inches(0.06), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN
        dot.line.fill.background()
    core = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, cx - Inches(0.5), cy - Inches(0.5), Inches(1), Inches(1))
    core.fill.solid()
    core.fill.fore_color.rgb = RGBColor(5, 50, 88)
    core.fill.transparency = 10
    core.line.color.rgb = CYAN
    set_text(core.text_frame, "AI", 24, WHITE, True, PP_ALIGN.CENTER)
    core.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def architecture(slide, x, y, w, h):
    layers = [
        ("感知层", ["IoT传感器", "数字孪生状态", "设备/物料/订单"]),
        ("决策层", ["11个AI智能体", "DynaSwarm拓扑", "ReflAct推理"]),
        ("执行层", ["CP-SAT排程", "甘特图可视化", "事件驱动重排"]),
    ]
    layer_h = h / 3 - Inches(0.11)
    for i, (name, items) in enumerate(layers):
        yy = y + i * (layer_h + Inches(0.17))
        glass_rect(slide, x, yy, w, layer_h, fill=PANEL_2, transparency=26)
        add_text(slide, name, x + Inches(0.24), yy + Inches(0.22), Inches(1.5), Inches(0.4), 18, CYAN, True)
        for j, item in enumerate(items):
            bx = x + Inches(2.05) + j * Inches(1.85)
            chip = glass_rect(slide, bx, yy + Inches(0.23), Inches(1.55), Inches(0.46), fill=RGBColor(6, 42, 76), transparency=18)
            set_text(chip.text_frame, item, 10, WHITE, True, PP_ALIGN.CENTER)
            chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i < 2:
            add_text(slide, "↓", x + w / 2 - Inches(0.1), yy + layer_h - Inches(0.02), Inches(0.3), Inches(0.3), 16, CYAN, True, PP_ALIGN.CENTER)


def topology_mini(slide, x, y, w, h, name, variant):
    glass_rect(slide, x, y, w, h, fill=RGBColor(7, 28, 58), transparency=18)
    add_text(slide, name, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2), Inches(0.22), 11, WHITE, True, PP_ALIGN.CENTER)
    nodes = []
    if variant == "chain":
        coords = [(0.35, 0.77), (0.85, 0.77), (1.35, 0.77)]
    elif variant == "fusion":
        coords = [(0.36, 0.62), (0.85, 0.92), (1.34, 0.62), (0.85, 0.45)]
    elif variant == "tree":
        coords = [(0.85, 0.45), (0.45, 0.9), (0.85, 0.9), (1.25, 0.9)]
    elif variant == "mesh":
        coords = [(0.45, 0.55), (1.15, 0.55), (0.45, 1.0), (1.15, 1.0)]
    else:
        coords = [(0.35, 0.55), (0.8, 0.55), (1.25, 0.55), (0.8, 1.0)]
    abs_nodes = [(x + Inches(a), y + Inches(b)) for a, b in coords]
    for a in range(len(abs_nodes) - 1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, abs_nodes[a][0] + Inches(0.05), abs_nodes[a][1] + Inches(0.05), abs_nodes[a + 1][0] + Inches(0.05), abs_nodes[a + 1][1] + Inches(0.05))
        line.line.color.rgb = CYAN
        line.line.transparency = 42
    for xx, yy in abs_nodes:
        n = slide.shapes.add_shape(MSO_SHAPE.OVAL, xx, yy, Inches(0.13), Inches(0.13))
        n.fill.solid()
        n.fill.fore_color.rgb = CYAN
        n.line.fill.background()


def huawei_stack(slide, x, y, w, h):
    services = [
        ("盘古大模型", "Agent推理"),
        ("MindIE", "昇腾推理"),
        ("GaussDB", "工业数据"),
        ("IoTDA", "设备接入"),
        ("EventGrid", "事件分发"),
        ("FunctionGraph", "事件触发"),
        ("ModelArts", "训练评估"),
    ]
    cols = 3
    cell_w = w / cols - Inches(0.12)
    cell_h = Inches(0.82)
    for i, (name, desc) in enumerate(services):
        row, col = divmod(i, cols)
        xx = x + col * (cell_w + Inches(0.18))
        yy = y + row * (cell_h + Inches(0.18))
        glass_rect(slide, xx, yy, cell_w, cell_h, fill=RGBColor(8, 36, 72), transparency=16)
        add_text(slide, name, xx + Inches(0.14), yy + Inches(0.14), cell_w - Inches(0.28), Inches(0.24), 13, WHITE, True)
        add_text(slide, desc, xx + Inches(0.14), yy + Inches(0.44), cell_w - Inches(0.28), Inches(0.2), 9, MUTED)


def dashboard_mock(slide, x, y, w, h, mode="normal"):
    glass_rect(slide, x, y, w, h, fill=RGBColor(5, 18, 43), transparency=8)
    add_text(slide, "DynaTwin-Swarm 工业数字孪生大屏", x + Inches(0.25), y + Inches(0.18), Inches(4.8), Inches(0.3), 15, WHITE, True)
    risk = "高风险审查" if mode == "fault" else ("恢复优化" if mode == "recover" else "实时调度")
    add_tag(slide, risk, x + w - Inches(1.65), y + Inches(0.16), Inches(1.28), RED if mode == "fault" else CYAN)
    # panels
    left = glass_rect(slide, x + Inches(0.22), y + Inches(0.68), w * 0.24, h - Inches(0.92), fill=RGBColor(8, 35, 70), transparency=18)
    mid = glass_rect(slide, x + Inches(2.65), y + Inches(0.68), w * 0.45, h * 0.54, fill=RGBColor(8, 35, 70), transparency=18)
    right = glass_rect(slide, x + w - Inches(2.75), y + Inches(0.68), Inches(2.48), h - Inches(0.92), fill=RGBColor(8, 35, 70), transparency=18)
    bottom = glass_rect(slide, x + Inches(2.65), y + h * 0.61, w * 0.45, h * 0.29, fill=RGBColor(8, 35, 70), transparency=18)
    for panel, label in [(left, "设备状态"), (mid, "3D工厂实时布局"), (right, "AI推理日志 / 订单队列"), (bottom, "排程甘特图")]:
        add_text(slide, label, panel.left + Inches(0.15), panel.top + Inches(0.13), panel.width - Inches(0.3), Inches(0.24), 10, CYAN, True)
    # machine bars
    for i in range(8):
        col, row = divmod(i, 4)
        xx = x + Inches(2.95) + col * Inches(0.78)
        yy = y + Inches(2.5) - row * Inches(0.38)
        height = Inches(0.35 + (i % 4) * 0.12)
        color = RED if mode == "fault" and i == 4 else (GREEN if mode == "recover" and i == 4 else CYAN)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xx, yy, Inches(0.35), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.fill.transparency = 10
        bar.line.fill.background()
    # machine list
    for i in range(6):
        yy = y + Inches(1.15 + i * 0.42)
        color = RED if mode == "fault" and i == 4 else (GREEN if mode == "recover" and i == 4 else CYAN)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.45), yy + Inches(0.06), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        temp = "107°C" if mode == "fault" and i == 4 else ("25°C" if mode == "recover" and i == 4 else f"{62+i*3}°C")
        add_text(slide, f"M{i+1}  {temp}", x + Inches(0.62), yy, Inches(1.55), Inches(0.2), 8, WHITE if i == 4 else MUTED, i == 4)
    # logs
    logs = ["Monitor: 检测设备状态", "Router: 选择拓扑", "Schedule: 重新求解", "Report: 输出方案"]
    if mode == "fault":
        logs = ["M5温度107°C告警", "切换高风险审查拓扑", "隔离故障机器", "CP-SAT重新排程"]
    if mode == "recover":
        logs = ["M5自动恢复", "重新纳入资源池", "优化排程方案", "完工时间缩短"]
    for i, log in enumerate(logs):
        add_text(slide, log, x + w - Inches(2.45), y + Inches(1.22 + i * 0.35), Inches(2.0), Inches(0.2), 8, MUTED)
    # gantt
    for i in range(5):
        yy = y + h * 0.68 + Inches(i * 0.21)
        add_text(slide, f"O{i+1}", x + Inches(2.9), yy, Inches(0.3), Inches(0.14), 6, DIM, font=MONO)
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(3.25), yy, Inches(3.55), Inches(0.1))
        track.fill.solid()
        track.fill.fore_color.rgb = RGBColor(25, 62, 96)
        track.line.fill.background()
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(3.25 + i * 0.28), yy, Inches(0.8 + i * 0.22), Inches(0.1))
        fill.fill.solid()
        fill.fill.fore_color.rgb = RED if mode == "fault" and i == 2 else CYAN
        fill.line.fill.background()


def table(slide, x, y, w, h, headers, rows):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    tbl = slide.shapes.add_table(rows_n, cols_n, x, y, w, h).table
    widths = [1.55, 1.15, 1.15, 1.7, 1.35, 1.1]
    for i, width in enumerate(widths[:cols_n]):
        tbl.columns[i].width = Inches(width)
    for c, head in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(7, 48, 83)
        cell.text = head
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(9, 29, 60) if r % 2 else RGBColor(12, 38, 72)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE if c == 0 else MUTED
            p.font.bold = c == 0
            p.alignment = PP_ALIGN.CENTER
    return tbl


def create_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    def new_slide():
        s = prs.slides.add_slide(blank)
        add_bg(s)
        return s

    # 1
    s = new_slide()
    add_text(s, "DynaTwin-Swarm", Inches(0.75), Inches(1.25), Inches(7.4), Inches(0.85), 48, WHITE, True)
    add_text(s, "工业数字孪生智能排产系统", Inches(0.78), Inches(2.2), Inches(7.2), Inches(0.55), 26, CYAN, True)
    add_text(s, "2026年贵州省人工智能创业大赛 · 昇腾AI创新赛题 · 贵州大学", Inches(0.8), Inches(6.55), Inches(8.7), Inches(0.28), 13, MUTED)
    add_multiline(s, ["多智能体协作 · 数字孪生 · CP-SAT排程 · 华为昇腾生态"], Inches(0.82), Inches(3.05), Inches(6.9), Inches(0.45), 17, MUTED)
    tech_decoration(s)

    # 2
    s = new_slide(); title(s, "目录", "六个章节快速讲清项目价值", 2)
    chapters = [("01", "问题背景"), ("02", "解决方案"), ("03", "核心技术"), ("04", "系统演示"), ("05", "数据验证"), ("06", "团队与规划")]
    for i, (no, name) in enumerate(chapters):
        col, row = i % 3, i // 3
        x, y = Inches(1.05 + col * 3.85), Inches(1.8 + row * 2.0)
        glass_rect(s, x, y, Inches(3.25), Inches(1.35), fill=RGBColor(6, 38, 78), transparency=12)
        add_text(s, no, x + Inches(0.28), y + Inches(0.24), Inches(0.55), Inches(0.32), 18, CYAN, True, font=MONO)
        add_text(s, name, x + Inches(0.28), y + Inches(0.66), Inches(2.4), Inches(0.36), 20, WHITE, True)

    # 3
    s = new_slide(); title(s, "传统工厂调度的三大痛点", "贵州本地工业智能化转型的现实需求", 3)
    card(s, Inches(0.75), Inches(1.75), Inches(3.75), Inches(2.45), "响应迟滞", "设备故障后人工重排需数小时\n每次停产损失数十万元", "⏱", CYAN)
    card(s, Inches(4.8), Inches(1.75), Inches(3.75), Inches(2.45), "人工依赖", "调度经验难以沉淀\n关键人员离职导致调度能力断层", "👤", AMBER)
    card(s, Inches(8.85), Inches(1.75), Inches(3.75), Inches(2.45), "扰动处理弱", "紧急插单、设备故障等动态扰动\n无法实时响应", "⚡", RED)
    glass_rect(s, Inches(1.1), Inches(5.0), Inches(11.1), Inches(0.85), fill=RGBColor(8, 42, 76), transparency=12)
    add_text(s, "贵州磷化工行业年产值超800亿，智能调度渗透率不足5%", Inches(1.38), Inches(5.25), Inches(10.5), Inches(0.3), 22, CYAN, True, PP_ALIGN.CENTER)

    # 4
    s = new_slide(); title(s, "DynaTwin-Swarm · 秒级响应的AI调度系统", "从感知、决策到执行的闭环工业AI", 4)
    architecture(s, Inches(0.9), Inches(1.65), Inches(7.1), Inches(4.55))
    glass_rect(s, Inches(8.45), Inches(1.65), Inches(3.95), Inches(4.55), fill=RGBColor(6, 38, 78), transparency=12)
    add_text(s, "核心价值", Inches(8.78), Inches(1.95), Inches(2.4), Inches(0.4), 22, CYAN, True)
    add_multiline(s, ["故障响应：小时级 → 秒级", "OEE提升：5-15%", "人工排产成本降低：70%"], Inches(8.82), Inches(2.75), Inches(3.05), Inches(1.7), 20, WHITE, gap=10)

    # 5
    s = new_slide(); title(s, "动态图结构自适应选择", "来源：EMNLP 2025 顶会论文 DynaSwarm", 5)
    variants = [("串行链", "chain"), ("并行汇聚", "fusion"), ("层级树", "tree"), ("全连接", "mesh"), ("高风险审查", "risk")]
    for i, (name, variant) in enumerate(variants):
        topology_mini(s, Inches(0.72 + i * 2.42), Inches(1.72), Inches(2.05), Inches(1.6), name, variant)
    glass_rect(s, Inches(1.0), Inches(4.05), Inches(11.25), Inches(1.0), fill=RGBColor(8, 42, 76), transparency=13)
    add_text(s, "基于A2C强化学习，根据工厂实时状态自动选择最优协作拓扑", Inches(1.35), Inches(4.27), Inches(10.6), Inches(0.35), 22, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, "不同场景用不同协作方式，高风险时多层审核保障安全", Inches(1.35), Inches(5.35), Inches(10.6), Inches(0.35), 22, CYAN, True, PP_ALIGN.CENTER)

    # 6
    s = new_slide(); title(s, "ReflAct目标状态反射推理", "来源：EMNLP 2025 顶会论文 ReflAct", 6)
    chain = ["观察当前状态", "对齐生产目标", "分析偏差", "输出建议动作"]
    for i, name in enumerate(chain):
        x = Inches(2.0 + i * 2.4)
        shp = glass_rect(s, x, Inches(2.55), Inches(1.75), Inches(0.8), fill=RGBColor(6, 45, 82), transparency=12)
        set_text(shp.text_frame, name, 13, WHITE, True, PP_ALIGN.CENTER)
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i < len(chain) - 1:
            line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(1.75), Inches(2.95), x + Inches(2.35), Inches(2.95))
            line.line.color.rgb = CYAN
            line.line.width = Pt(1.6)
    card(s, Inches(0.85), Inches(4.35), Inches(5.45), Inches(1.25), "智能体协作", "11个专业智能体分工协作，豆包大模型驱动推理", None, CYAN)
    card(s, Inches(7.0), Inches(4.35), Inches(5.45), Inches(1.25), "可解释可审计", "输出观察、目标、偏差和建议动作，避免不可控黑箱决策", None, VIOLET)

    # 7
    s = new_slide(); title(s, "深度集成华为昇腾AI技术体系", "华为云 + 昇腾AI + 工业数字孪生", 7)
    huawei_stack(s, Inches(1.1), Inches(1.75), Inches(7.3), Inches(3.5))
    glass_rect(s, Inches(9.0), Inches(1.75), Inches(3.2), Inches(3.5), fill=RGBColor(6, 38, 78), transparency=12)
    add_text(s, "部署价值", Inches(9.3), Inches(2.05), Inches(1.8), Inches(0.35), 22, CYAN, True)
    add_multiline(s, ["全栈接入华为云", "西南-贵阳一区域部署", "体现昇腾算力价值", "支持Mock到真实服务平滑切换"], Inches(9.32), Inches(2.75), Inches(2.45), Inches(1.75), 16, WHITE, gap=5)

    # 8
    s = new_slide(); title(s, "一屏掌控全局 · 工业数字孪生大屏", "自动运行模式下工厂状态实时刷新，故障秒级感知自动重排", 8)
    dashboard_mock(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(4.95), "normal")

    # 9
    s = new_slide(); title(s, "场景演示：设备故障自动重排", "故障感知、拓扑切换、智能体协作、CP-SAT重排", 9)
    add_multiline(s, ["① M5机器温度107°C触发告警", "② AI自动切换高风险审查拓扑", "③ 11个智能体协作分析", "④ CP-SAT重新求解排程", "⑤ 甘特图实时更新"], Inches(0.85), Inches(1.75), Inches(4.1), Inches(3.2), 18, WHITE, gap=8)
    dashboard_mock(s, Inches(5.0), Inches(1.45), Inches(7.45), Inches(4.95), "fault")

    # 10
    s = new_slide(); title(s, "场景演示：机器自动恢复", "恢复事件触发后重新纳入调度资源池", 10)
    glass_rect(s, Inches(0.85), Inches(1.8), Inches(4.0), Inches(3.25), fill=RGBColor(8, 42, 76), transparency=12)
    add_multiline(s, ["故障机器经过随机30-120分钟后自动恢复", "AI感知恢复事件", "重新纳入调度资源池", "优化排程方案", "完工时间缩短"], Inches(1.15), Inches(2.15), Inches(3.4), Inches(2.3), 18, WHITE, gap=7)
    dashboard_mock(s, Inches(5.0), Inches(1.45), Inches(7.45), Inches(4.95), "recover")

    # 11
    s = new_slide(); title(s, "5个国际标准数据集验证系统性能", "求解时间均在200ms以内，验证系统实时响应能力", 11)
    headers = ["数据集名称", "作业数", "机器数", "本系统Makespan", "最优已知值", "差距%"]
    rows = [
        ["FT06", 6, 6, 67, 55, "21.8%"],
        ["ABZ9", 20, 15, 873, 678, "28.8%"],
        ["LA40", 15, 15, 1307, 1222, "7.0%"],
        ["SWV20", 50, 10, 2823, 1675, "68.5%"],
        ["DMU80", 50, 20, 9280, "暂无", "-"],
    ]
    table(s, Inches(0.75), Inches(1.65), Inches(11.85), Inches(3.85), headers, rows)
    add_text(s, "核心结论：系统已具备公共Benchmark评测闭环，可持续优化拓扑选择与调度求解质量。", Inches(1.0), Inches(6.0), Inches(11.2), Inches(0.35), 17, CYAN, True, PP_ALIGN.CENTER)

    # 12
    s = new_slide(); title(s, "贵州重点工业场景深度适配", "围绕贵州本地优势产业和工业互联网场景落地", 12)
    scenarios = [
        ("磷化工生产调度", "贵州最大磷矿基地，年产值800亿+"),
        ("非煤矿山作业调度", "高危环境安全保障"),
        ("智能制造产线优化", "万企融合战略"),
        ("工业园区数字孪生", "IoT全链路贯通"),
    ]
    for i, (head, body) in enumerate(scenarios):
        card(s, Inches(0.8 + i * 3.1), Inches(2.0), Inches(2.65), Inches(2.55), head, body, None, CYAN if i % 2 == 0 else VIOLET)

    # 13
    s = new_slide(); title(s, "可持续的商业化路径", "从SaaS订阅到行业知识库沉淀", 13)
    paths = [
        ("短期", "SaaS订阅\n面向中小型制造企业\n按月收费"),
        ("中期", "定制集成\n面向大型工业企业\n项目制交付"),
        ("长期", "数据智能\n积累调度知识库\n输出行业解决方案"),
    ]
    for i, (head, body) in enumerate(paths):
        card(s, Inches(0.95 + i * 4.05), Inches(1.85), Inches(3.45), Inches(2.75), head, body, None, [CYAN, AMBER, GREEN][i])
    glass_rect(s, Inches(1.25), Inches(5.25), Inches(10.8), Inches(0.65), fill=RGBColor(8, 42, 76), transparency=12)
    add_text(s, "目标市场：贵州及西南地区制造业企业，市场规模超500亿", Inches(1.4), Inches(5.43), Inches(10.5), Inches(0.26), 19, CYAN, True, PP_ALIGN.CENTER)

    # 14
    s = new_slide(); title(s, "团队背景", "贵州大学 · 贵州省公共大数据重点实验室", 14)
    glass_rect(s, Inches(0.95), Inches(1.75), Inches(4.2), Inches(3.8), fill=RGBColor(6, 38, 78), transparency=12)
    add_text(s, "王骏", Inches(1.25), Inches(2.1), Inches(3.4), Inches(0.52), 32, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, "贵州大学网络与信息安全专业硕士研究生", Inches(1.25), Inches(2.78), Inches(3.4), Inches(0.3), 15, CYAN, True, PP_ALIGN.CENTER)
    add_text(s, "指导教师：谭伟杰", Inches(1.25), Inches(3.35), Inches(3.4), Inches(0.28), 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, "GitHub：github.com/wjhard/DynaTwin-Swarm", Inches(1.25), Inches(4.65), Inches(3.4), Inches(0.3), 11, MUTED, False, PP_ALIGN.CENTER)
    add_multiline(s, ["平台：贵州省公共大数据重点实验室", "研究方向：IoT安全与多智能体系统", "工程能力：工业数字孪生、AI Agent、调度优化、可视化大屏"], Inches(6.0), Inches(2.15), Inches(5.6), Inches(2.0), 20, WHITE, gap=9)

    # 15
    s = new_slide(); title(s, "核心价值总结", "感谢评委老师 · 期待交流指导", 15)
    card(s, Inches(0.75), Inches(1.75), Inches(3.85), Inches(2.45), "顶会论文工程化落地", "EMNLP 2025两篇论文实现", "📄", CYAN)
    card(s, Inches(4.75), Inches(1.75), Inches(3.85), Inches(2.45), "华为昇腾全栈集成", "深度对接贵阳一区域算力", "🔷", VIOLET)
    card(s, Inches(8.75), Inches(1.75), Inches(3.85), Inches(2.45), "贵州工业场景适配", "磷化工矿山落地方案", "🏭", GREEN)
    add_text(s, "感谢评委老师 · 期待交流指导", Inches(1.3), Inches(5.35), Inches(10.8), Inches(0.45), 28, CYAN, True, PP_ALIGN.CENTER)
    add_text(s, "联系方式：贵州大学网络与信息安全专业硕士研究生团队 · GitHub: github.com/wjhard/DynaTwin-Swarm", Inches(1.35), Inches(6.08), Inches(10.7), Inches(0.28), 13, MUTED, False, PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    create_deck()
